"""Generate Salanor Aegis pitch deck (.pptx). Run: python docs/build-aegis-deck.py"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
LOGO = ROOT / "apps" / "web-marketing" / "public" / "salanor-logo.png"
OUT = Path.home() / "Downloads" / "Salanor_Aegis_Pitch_Deck.pptx"

# Brand palette (matches marketing site)
BG = RGBColor(0x0A, 0x0C, 0x10)
SURFACE = RGBColor(0x12, 0x14, 0x1A)
BORDER = RGBColor(0x27, 0x27, 0x2A)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
TEAL_DIM = RGBColor(0x14, 0x5A, 0x50)
WHITE = RGBColor(0xF4, 0xF4, 0xF5)
MUTED = RGBColor(0xA1, 0xA1, 0xAA)
DIM = RGBColor(0x71, 0x71, 0x7A)

FONT = "Segoe UI"
MARGIN_L = Inches(0.65)
MARGIN_R = Inches(0.65)
CONTENT_W = Inches(8.7)


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_round_card(slide, left, top, width, height) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    return shape


def write_text(
    shape_or_slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = WHITE,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
) -> None:
    if hasattr(shape_or_slide, "text_frame"):
        box = shape_or_slide
    else:
        box = shape_or_slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def write_lines(
    slide,
    left,
    top,
    width,
    height,
    lines: list[tuple[str, int, RGBColor, bool]],
    gap: int = 5,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.space_after = Pt(gap)
        if not text:
            continue
        run = p.runs[0] if p.runs else p.add_run()
        if not p.runs:
            run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT


def slide_header(slide, title: str, subtitle: str | None = None) -> None:
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), MARGIN_L, Inches(0.42), width=Inches(0.38))
    write_text(
        slide,
        Inches(1.08),
        Inches(0.48),
        Inches(1.2),
        Inches(0.25),
        "SALANOR",
        size=9,
        color=DIM,
        bold=True,
    )
    write_text(slide, MARGIN_L, Inches(0.95), CONTENT_W, Inches(0.65), title, size=24, bold=True)
    add_rect(slide, MARGIN_L, Inches(1.58), Inches(1.1), Inches(0.045), TEAL)
    if subtitle:
        write_text(slide, MARGIN_L, Inches(1.72), CONTENT_W, Inches(0.45), subtitle, size=13, color=MUTED)


def slide_footer(slide, text: str) -> None:
    write_text(slide, MARGIN_L, Inches(7.05), CONTENT_W, Inches(0.3), text, size=9, color=DIM)


def title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), TEAL_DIM)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), MARGIN_L, Inches(1.35), width=Inches(1.05))
    write_text(slide, MARGIN_L, Inches(2.55), Inches(3), Inches(0.35), "SALANOR", size=12, color=DIM, bold=True)
    write_text(slide, MARGIN_L, Inches(2.95), Inches(6), Inches(0.9), "Aegis", size=52, bold=True, color=TEAL)
    write_text(
        slide,
        MARGIN_L,
        Inches(4.05),
        Inches(8.2),
        Inches(0.9),
        "Proof of what your agents actually did.",
        size=30,
        bold=True,
    )
    write_text(
        slide,
        MARGIN_L,
        Inches(4.95),
        Inches(8.2),
        Inches(0.85),
        "Signed, tamper-evident records for teams who have to answer when something goes wrong.",
        size=15,
        color=MUTED,
    )
    slide_footer(slide, "Salanor Ltd  |  Kigali, Rwanda  |  Design partners  |  2026")


def hook_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    write_text(slide, MARGIN_L, Inches(2.0), Inches(8.5), Inches(1.2), "Nobody can explain\nwhat the agent did.", size=36, bold=True)
    write_text(
        slide,
        MARGIN_L,
        Inches(3.55),
        Inches(8.0),
        Inches(1.2),
        "Teams want agents in loans, triage, and ops. The blocker is not the model.\nIt is what happens after something goes wrong.",
        size=16,
        color=MUTED,
    )
    add_rect(slide, MARGIN_L, Inches(5.0), Inches(2.5), Inches(0.04), TEAL)


def gaps_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "Three gaps we keep seeing")
    cards = [
        ("No record", "Most teams cannot reconstruct the sequence after a workflow runs. Logs are incomplete or scattered."),
        ("No proof", "What exists can be edited after the fact. That is not enough for a regulator, auditor, or counsel."),
        ("No owner", "Legal, security, and engineering each assume someone else is tracking it. Often nobody is."),
    ]
    w = Inches(2.75)
    gap = Inches(0.22)
    x0 = MARGIN_L
    y = Inches(2.05)
    h = Inches(3.35)
    for i, (head, body) in enumerate(cards):
        x = x0 + i * (w + gap)
        add_round_card(slide, x, y, w, h)
        write_text(slide, x + Inches(0.2), y + Inches(0.22), Inches(0.4), Inches(0.35), str(i + 1), size=18, bold=True, color=TEAL)
        write_text(slide, x + Inches(0.2), y + Inches(0.65), w - Inches(0.35), Inches(0.45), head, size=15, bold=True)
        write_text(slide, x + Inches(0.2), y + Inches(1.15), w - Inches(0.35), Inches(2.0), body, size=12, color=MUTED)


def quote_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_round_card(slide, Inches(1.2), Inches(2.0), Inches(7.6), Inches(2.8))
    write_text(
        slide,
        Inches(1.55),
        Inches(2.45),
        Inches(6.9),
        Inches(1.6),
        '"The gap is not model safety. It is proving what happened on a Tuesday in March, six months later."',
        size=20,
        color=TEAL,
        bold=True,
    )
    write_text(
        slide,
        Inches(1.55),
        Inches(4.15),
        Inches(6.9),
        Inches(0.4),
        "Model risk lead, design partner (top-15 US bank)",
        size=12,
        color=MUTED,
    )


def why_now_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "The infrastructure did not keep up")
    col_w = Inches(4.15)
    y = Inches(2.05)
    h = Inches(3.8)
    add_round_card(slide, MARGIN_L, y, col_w, h)
    add_round_card(slide, MARGIN_L + col_w + Inches(0.35), y, col_w, h)
    write_text(slide, MARGIN_L + Inches(0.22), y + Inches(0.2), col_w, Inches(0.35), "Agents shipped fast", size=14, bold=True)
    write_lines(
        slide,
        MARGIN_L + Inches(0.22),
        y + Inches(0.6),
        col_w - Inches(0.35),
        Inches(2.8),
        [
            ("LangGraph, CrewAI, OpenAI Agents SDK, MCP.", 12, MUTED, False),
            ("", 6, MUTED, False),
            ("Built for capability. Not for evidence.", 12, WHITE, False),
        ],
    )
    x2 = MARGIN_L + col_w + Inches(0.35)
    write_text(slide, x2 + Inches(0.22), y + Inches(0.2), col_w, Inches(0.35), "Rules caught up", size=14, bold=True)
    write_lines(
        slide,
        x2 + Inches(0.22),
        y + Inches(0.6),
        col_w - Inches(0.35),
        Inches(2.8),
        [
            ("EU AI Act (Arts. 12, 14): logging and oversight.", 12, MUTED, False),
            ("SOC 2 auditors ask agent-specific questions.", 12, MUTED, False),
            ("NIST AI RMF expects documented provenance.", 12, MUTED, False),
        ],
    )
    write_text(
        slide,
        MARGIN_L,
        Inches(6.15),
        CONTENT_W,
        Inches(0.55),
        "We did not start Salanor because agents are exciting. We started it because someone has to build the receipts before incidents pile up.",
        size=12,
        color=MUTED,
    )


def solution_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(
        slide,
        "One integration. A signed record of everything.",
        "Connect via SDK, n8n node, or API. No rewrite of your agent logic.",
    )
    steps = [
        ("Capture", "Every tool call and LLM turn, as it happens."),
        ("Sign", "Ed25519. Keys stay in your KMS."),
        ("Enforce", "Policy before execution. Under 5ms at p50."),
        ("Chain", "Append-only ledger. Public Merkle roots."),
        ("Export", "Compliance bundle when audit asks."),
    ]
    w = Inches(1.62)
    gap = Inches(0.14)
    y = Inches(2.35)
    h = Inches(2.55)
    for i, (head, body) in enumerate(steps):
        x = MARGIN_L + i * (w + gap)
        add_round_card(slide, x, y, w, h)
        write_text(slide, x + Inches(0.12), y + Inches(0.15), Inches(0.35), Inches(0.3), str(i + 1), size=14, bold=True, color=TEAL)
        write_text(slide, x + Inches(0.12), y + Inches(0.48), w - Inches(0.2), Inches(0.35), head, size=12, bold=True)
        write_text(slide, x + Inches(0.12), y + Inches(0.88), w - Inches(0.2), Inches(1.4), body, size=10, color=MUTED)
    write_text(
        slide,
        MARGIN_L,
        Inches(5.2),
        CONTENT_W,
        Inches(0.4),
        "Integrations: TypeScript, Python, Go SDKs  |  Salanor Aegis n8n node  |  REST API",
        size=11,
        color=DIM,
    )


def flow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "From agent action to audit-ready proof")
    steps = [
        ("Instrument", "SDK, n8n, or API on LangGraph, CrewAI, OpenAI Agents, MCP."),
        ("Sign & enforce", "Signed at the point of action. Keys never leave your infra."),
        ("Ledger", "Hash-chained storage. Third parties verify without your keys."),
        ("Audit", "Reconstruct a trace in ~30s. OTel to Splunk, Datadog, Sentinel."),
    ]
    y0 = 2.15
    for i, (head, body) in enumerate(steps):
        top = Inches(y0 + i * 1.15)
        add_round_card(slide, MARGIN_L, top, CONTENT_W, Inches(0.95))
        write_text(slide, MARGIN_L + Inches(0.18), top + Inches(0.12), Inches(0.35), Inches(0.3), str(i + 1), size=13, bold=True, color=TEAL)
        write_text(slide, MARGIN_L + Inches(0.55), top + Inches(0.1), Inches(2.2), Inches(0.3), head, size=13, bold=True)
        write_text(slide, MARGIN_L + Inches(2.85), top + Inches(0.12), Inches(5.5), Inches(0.65), body, size=11, color=MUTED)


def console_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "What the console gives you")
    add_round_card(slide, MARGIN_L, Inches(2.05), Inches(4.0), Inches(3.55))
    write_text(
        slide,
        MARGIN_L + Inches(0.35),
        Inches(3.35),
        Inches(3.3),
        Inches(0.9),
        "[ Drop console screenshot here ]",
        size=13,
        color=DIM,
        align=PP_ALIGN.CENTER,
    )
    items = [
        ("Trace view", "Every action in order, signed proof on each step."),
        ("Policy view", "Allowed, blocked, and why, before it runs."),
        ("Compliance export", "Built for the auditor, not a raw dump."),
        ("Public verify", "Check a trace was not altered, without your keys."),
    ]
    x = Inches(4.95)
    for i, (head, body) in enumerate(items):
        y = Inches(2.05 + i * 0.88)
        write_text(slide, x, y, Inches(3.8), Inches(0.28), head, size=13, bold=True)
        write_text(slide, x, y + Inches(0.28), Inches(3.8), Inches(0.5), body, size=11, color=MUTED)
    slide_footer(slide, "docs.salanor.com  |  app.salanor.com")


def moat_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "The format is open. The infrastructure is what we sell.")
    rows = [
        ("APS-1", "Open provenance standard for agent events (CC BY 4.0)."),
        ("Verifier CLI", "MIT-licensed. Audit a trace without Salanor online."),
        ("Control plane", "Ingest, policy, storage, exports. That is what we sell."),
        ("BYOK", "Ed25519 in your KMS. We cannot rewrite your history."),
    ]
    y = 2.05
    for head, body in rows:
        add_round_card(slide, MARGIN_L, Inches(y), CONTENT_W, Inches(0.88))
        write_text(slide, MARGIN_L + Inches(0.2), Inches(y + 0.12), Inches(1.5), Inches(0.3), head, size=13, bold=True, color=TEAL)
        write_text(slide, MARGIN_L + Inches(1.85), Inches(y + 0.15), Inches(6.5), Inches(0.55), body, size=12, color=MUTED)
        y += 1.02
    write_text(
        slide,
        MARGIN_L,
        Inches(6.2),
        CONTENT_W,
        Inches(0.4),
        "Logs and SIEM record activity. They do not produce signed, policy-bound, audit-ready provenance.",
        size=11,
        color=DIM,
    )


def compliance_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "Built toward frameworks your auditors already use")
    left = [
        "SOC 2 Type II  |  target Q4 2026",
        "EU AI Act  |  Arts. 12, 14, 19, 26",
        "NIST AI RMF  |  mapped exports",
        "HIPAA  |  BYOC and on-prem",
        "FedRAMP Moderate  |  target Q2 2027",
    ]
    write_lines(
        slide,
        MARGIN_L,
        Inches(2.05),
        Inches(4.2),
        Inches(3.0),
        [(line, 12, MUTED, False) for line in left],
        gap=8,
    )
    add_round_card(slide, Inches(5.15), Inches(2.05), Inches(3.95), Inches(2.35))
    write_text(slide, Inches(5.35), Inches(2.25), Inches(3.5), Inches(0.3), "Integrations", size=13, bold=True)
    write_lines(
        slide,
        Inches(5.35),
        Inches(2.65),
        Inches(3.5),
        Inches(1.5),
        [
            ("Keys: AWS KMS, GCP KMS, Vault", 11, MUTED, False),
            ("Export: Splunk, Datadog, Sentinel", 11, MUTED, False),
        ],
    )
    write_text(
        slide,
        MARGIN_L,
        Inches(5.35),
        CONTENT_W,
        Inches(0.65),
        "Admissibility depends on jurisdiction and counsel. Aegis supports evidentiary workflows; it does not guarantee court outcomes.",
        size=10,
        color=DIM,
    )


def roadmap_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "Aegis now. Aether next.")
    cards = [
        ("Aegis", "Provenance and audit. GA target Q4 2026.", "Now"),
        ("APS-1", "Open standard and verifier. Shipped.", "Live"),
        ("Aether", "Risk scoring and insurer telemetry. Research 2027.", "Next"),
    ]
    w = Inches(2.75)
    y = Inches(2.15)
    h = Inches(2.5)
    for i, (name, body, tag) in enumerate(cards):
        x = MARGIN_L + i * (w + Inches(0.22))
        add_round_card(slide, x, y, w, h)
        write_text(slide, x + Inches(0.18), y + Inches(0.18), Inches(1.2), Inches(0.25), tag.upper(), size=9, color=TEAL, bold=True)
        write_text(slide, x + Inches(0.18), y + Inches(0.48), w, Inches(0.35), name, size=16, bold=True)
        write_text(slide, x + Inches(0.18), y + Inches(0.95), w - Inches(0.3), Inches(1.3), body, size=11, color=MUTED)
    write_text(
        slide,
        MARGIN_L,
        Inches(5.0),
        CONTENT_W,
        Inches(1.0),
        "2027: reinsurers and carriers as Aether design partners on risk scoring. Not the Aegis GA buyer.\nAether reads signed patterns from Aegis, not raw payloads.",
        size=12,
        color=MUTED,
    )


def icp_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "Who actually needs this", "Today's buyer is governance and compliance, not insurers.")
    cards = [
        ("Banks & fintech", "Model risk before agents touch underwriting or fraud."),
        ("Health systems", "No autonomous triage without a reviewable record."),
        ("Enterprise platforms", "Customers who need proof of what the product did."),
    ]
    w = Inches(2.75)
    y = Inches(2.35)
    h = Inches(2.65)
    for i, (head, body) in enumerate(cards):
        x = MARGIN_L + i * (w + Inches(0.22))
        add_round_card(slide, x, y, w, h)
        write_text(slide, x + Inches(0.18), y + Inches(0.22), w - Inches(0.3), Inches(0.55), head, size=14, bold=True)
        write_text(slide, x + Inches(0.18), y + Inches(0.85), w - Inches(0.3), Inches(1.5), body, size=11, color=MUTED)
    write_text(
        slide,
        MARGIN_L,
        Inches(5.35),
        CONTENT_W,
        Inches(0.55),
        "Buyer: AI governance, security, or compliance.  Champion: the engineer who would otherwise build the audit trail alone.",
        size=12,
        color=MUTED,
    )


def traction_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "Where we are")
    metrics = [
        ("Stage", "Prototype. Design partner pilots."),
        ("HQ", "Kigali, Rwanda (Salanor Ltd). Founded 2020."),
        ("GTM", "Sales-led. Direct onboarding."),
    ]
    for i, (label, val) in enumerate(metrics):
        y = 2.05 + i * 0.95
        top = Inches(y)
        add_round_card(slide, MARGIN_L, top, CONTENT_W, Inches(0.78))
        write_text(slide, MARGIN_L + Inches(0.2), top + Inches(0.14), Inches(1.2), Inches(0.3), label, size=12, bold=True, color=TEAL)
        write_text(slide, MARGIN_L + Inches(1.5), top + Inches(0.16), Inches(6.8), Inches(0.45), val, size=12, color=MUTED)
    write_text(slide, MARGIN_L, Inches(5.05), CONTENT_W, Inches(0.3), "In active development", size=13, bold=True)
    write_lines(
        slide,
        MARGIN_L,
        Inches(5.4),
        CONTENT_W,
        Inches(1.2),
        [
            ("Console (traces, policy, exports)", 11, MUTED, False),
            ("APS-1 ingest and policy enforcement", 11, MUTED, False),
            ("Compliance export bundles", 11, MUTED, False),
            ("n8n node and SDK integrations", 11, MUTED, False),
        ],
        gap=3,
    )


def business_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "How we make money")
    add_round_card(slide, MARGIN_L, Inches(2.05), Inches(4.15), Inches(2.2))
    add_round_card(slide, Inches(5.15), Inches(2.05), Inches(3.95), Inches(2.2))
    write_text(slide, MARGIN_L + Inches(0.2), Inches(2.25), Inches(3.5), Inches(0.3), "SaaS control plane", size=14, bold=True)
    write_text(
        slide,
        MARGIN_L + Inches(0.2),
        Inches(2.65),
        Inches(3.7),
        Inches(1.4),
        "Per organization, tiered by event volume and team size.",
        size=12,
        color=MUTED,
    )
    write_text(slide, Inches(5.35), Inches(2.25), Inches(3.5), Inches(0.3), "Enterprise", size=14, bold=True)
    write_text(
        slide,
        Inches(5.35),
        Inches(2.65),
        Inches(3.5),
        Inches(1.4),
        "SSO, dedicated region, custom retention, SLA.",
        size=12,
        color=MUTED,
    )
    write_text(
        slide,
        MARGIN_L,
        Inches(4.55),
        CONTENT_W,
        Inches(0.9),
        "APS-1 stays free. Revenue is running Aegis, not owning the format.\nNo self-serve billing yet. Early customers are ops-led while we get pricing right.",
        size=12,
        color=MUTED,
    )


def team_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "Team")
    add_round_card(slide, MARGIN_L, Inches(2.05), CONTENT_W, Inches(2.85))
    write_text(slide, MARGIN_L + Inches(0.25), Inches(2.3), CONTENT_W, Inches(0.35), "Landry Bougang Fotso", size=16, bold=True)
    write_text(slide, MARGIN_L + Inches(0.25), Inches(2.65), CONTENT_W, Inches(0.25), "Founder & CEO", size=12, color=TEAL)
    write_text(
        slide,
        MARGIN_L + Inches(0.25),
        Inches(3.05),
        Inches(8.0),
        Inches(0.9),
        "Building infrastructure so high-stakes AI decisions can be explained after the fact, not only made in the moment.",
        size=12,
        color=MUTED,
    )
    write_text(
        slide,
        MARGIN_L + Inches(0.25),
        Inches(3.95),
        Inches(8.0),
        Inches(0.7),
        "Earlier work in environmental sensing in Rwanda. Pivoted to Aegis after customer discovery pointed at a sharper problem.",
        size=11,
        color=DIM,
    )


def ask_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "Design partners")
    checks = [
        "You run, or are about to run, agents where being wrong is expensive.",
        "You need signed provenance before scaling past a pilot.",
        "Compliance or legal will help shape exports they would actually accept.",
    ]
    for i, line in enumerate(checks):
        y = Inches(2.15 + i * 0.72)
        write_text(slide, MARGIN_L, y, Inches(0.35), Inches(0.3), "✓", size=14, color=TEAL, bold=True)
        write_text(slide, MARGIN_L + Inches(0.35), y, Inches(8.0), Inches(0.55), line, size=13, color=MUTED)
    add_round_card(slide, MARGIN_L, Inches(4.55), CONTENT_W, Inches(1.05))
    write_text(slide, MARGIN_L + Inches(0.25), Inches(4.78), CONTENT_W, Inches(0.35), "partners@salanor.com", size=14, bold=True, color=TEAL)
    write_text(slide, MARGIN_L + Inches(0.25), Inches(5.15), CONTENT_W, Inches(0.3), "https://www.salanor.com/contact", size=12, color=MUTED)


def close_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(7.44), Inches(10), Inches(0.06), TEAL_DIM)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.55), Inches(1.85), width=Inches(0.95))
    write_text(slide, Inches(0.7), Inches(3.05), Inches(8.6), Inches(0.4), "SALANOR", size=12, color=DIM, align=PP_ALIGN.CENTER, bold=True)
    write_text(
        slide,
        Inches(0.7),
        Inches(3.45),
        Inches(8.6),
        Inches(0.7),
        "Aegis: provenance you can take to audit.",
        size=26,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    write_text(
        slide,
        Inches(0.7),
        Inches(4.35),
        Inches(8.6),
        Inches(0.4),
        "www.salanor.com  |  docs.salanor.com  |  salanor.com/spec",
        size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    hook_slide(prs)
    gaps_slide(prs)
    quote_slide(prs)
    why_now_slide(prs)
    solution_slide(prs)
    flow_slide(prs)
    console_slide(prs)
    moat_slide(prs)
    compliance_slide(prs)
    roadmap_slide(prs)
    icp_slide(prs)
    traction_slide(prs)
    business_slide(prs)
    team_slide(prs)
    ask_slide(prs)
    close_slide(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    print(f"Wrote {build()}")
