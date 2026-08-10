"""Aegis design-partner pitch deck — human voice, sourced proof, connectors."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(r"c:\Users\landry\Downloads\Aegis_Pitch_Deck_Salanor.pptx")
OUT_REPO = Path(r"D:\PROJECTS\salanor\docs\Aegis_Pitch_Deck.pptx")

BG = RGBColor(0x0B, 0x0F, 0x14)
CARD = RGBColor(0x14, 0x1A, 0x22)
LINE = RGBColor(0x2A, 0x33, 0x40)
TEXT = RGBColor(0xE8, 0xEC, 0xF1)
MUTED = RGBColor(0x9A, 0xA5, 0xB5)
ACCENT = RGBColor(0xC8, 0xA4, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W = Inches(13.333)
H = Inches(7.5)
TOTAL = 13


def _set_run(run, size=18, bold=False, color=TEXT, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    _fill(bg, BG)
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return slide


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color)
    return box


def add_paras(slide, left, top, width, height, lines, size=15, color=TEXT, spacing=8):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        text, bold = (item if isinstance(item, tuple) else (item, False))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = text
        _set_run(run, size=size, bold=bold, color=color)
    return box


def footer(slide, n):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3), "Salanor  ·  Aegis", size=11, color=MUTED)
    add_text(slide, Inches(11.8), Inches(7.05), Inches(1), Inches(0.3), str(n), size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def accent_bar(slide, left, top):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), Inches(0.45))
    _fill(bar, ACCENT)


def card(slide, left, top, width, height):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill(sh, CARD)
    sh.line.color.rgb = LINE
    sh.line.width = Pt(1)
    return sh


def title_block(slide, title):
    accent_bar(slide, Inches(0.8), Inches(0.55))
    add_text(slide, Inches(1.05), Inches(0.5), Inches(11.5), Inches(0.5), title, size=26, bold=True, color=WHITE)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 Title
    s = blank_slide(prs)
    add_text(s, Inches(0.8), Inches(2.0), Inches(11), Inches(0.35), "SALANOR", size=14, bold=True, color=ACCENT)
    add_text(s, Inches(0.8), Inches(2.45), Inches(11.5), Inches(0.85), "Aegis", size=52, bold=True, color=WHITE)
    add_text(
        s,
        Inches(0.8),
        Inches(3.4),
        Inches(11),
        Inches(0.7),
        "A signed record of what your AI agents actually did.",
        size=24,
        color=TEXT,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(4.3),
        Inches(10.5),
        Inches(0.9),
        "For compliance, security, and the engineer who gets the message\nwhen something goes wrong.",
        size=16,
        color=MUTED,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(6.2),
        Inches(11),
        Inches(0.35),
        "Salanor Ltd  ·  Kigali  ·  Design partner program  ·  2026",
        size=13,
        color=MUTED,
    )

    # 2 Problem
    s = blank_slide(prs)
    title_block(s, "The gap isn't model safety")
    add_text(
        s,
        Inches(0.8),
        Inches(1.15),
        Inches(11.5),
        Inches(0.6),
        "Teams want agents on real work — loans, triage, ops. The pain shows up after a bad decision:\nnobody can reconstruct what happened in a form an auditor will accept.",
        size=15,
        color=MUTED,
    )
    items = [
        ("No record", "Dozens of steps. Logs missing, incomplete, or split across three tools."),
        ("No proof", "Whatever exists can be edited later. That doesn't hold with a regulator."),
        ("No owner", "Legal, security, and engineering each assume someone else is watching."),
    ]
    for i, (t, b) in enumerate(items):
        x = Inches(0.8) + Inches(4.0) * i
        card(s, x, Inches(2.1), Inches(3.7), Inches(2.5))
        add_text(s, x + Inches(0.25), Inches(2.3), Inches(3.2), Inches(0.3), str(i + 1), size=13, bold=True, color=ACCENT)
        add_text(s, x + Inches(0.25), Inches(2.7), Inches(3.2), Inches(0.4), t, size=18, bold=True, color=WHITE)
        add_text(s, x + Inches(0.25), Inches(3.25), Inches(3.2), Inches(1.1), b, size=14, color=MUTED)
    add_text(
        s,
        Inches(0.8),
        Inches(5.1),
        Inches(11.5),
        Inches(1.0),
        '"The gap isn\'t model safety. It\'s proving what happened on a Tuesday in March, six months later."\n— Model risk lead, top-15 US bank (design partner conversation)',
        size=14,
        color=TEXT,
    )
    footer(s, 2)

    # 3 Proof — sourced, human
    s = blank_slide(prs)
    title_block(s, "Why this is a real problem — not a vibe")
    add_text(
        s,
        Inches(0.8),
        Inches(1.15),
        Inches(11.5),
        Inches(0.4),
        "We don't invent a $47B scare number. Here's what you can actually point to:",
        size=15,
        color=MUTED,
    )
    proofs = [
        (
            "The law already asks for logs",
            "EU AI Act Art. 12: high-risk systems must automatically record events for traceability.\nArt. 26: deployers keep those logs (at least 6 months). Miss obligations → fines up to\n€15M or 3% of worldwide turnover (Art. 99). Prohibited practices go higher (up to 7%).",
        ),
        (
            "\"The bot did it\" already lost in court",
            "Moffatt v Air Canada (2024, BC tribunal): airline held liable for chatbot misinformation.\nTribunal rejected the idea that the chatbot is a separate entity. You own what it said.",
        ),
        (
            "Insurers are drawing lines around AI",
            "From 2026, ISO generative-AI exclusion endorsements started attaching to many US\nCGL renewals. Coverage for agent mistakes is no longer assumed — evidence matters.",
        ),
        (
            "Auditors are already asking",
            "When AI is in SOC 2 scope, reviewers want attributable trails: who called what model,\nwith what tools, and whether that log can be trusted (not rewritten by the app itself).",
        ),
    ]
    for i, (t, b) in enumerate(proofs):
        y = Inches(1.6) + Inches(1.2) * i
        add_text(s, Inches(0.8), y, Inches(11.5), Inches(0.3), t, size=15, bold=True, color=ACCENT)
        add_text(s, Inches(0.8), y + Inches(0.32), Inches(11.7), Inches(0.75), b, size=13, color=TEXT)
    add_text(
        s,
        Inches(0.8),
        Inches(6.5),
        Inches(11.5),
        Inches(0.35),
        "Sources: Regulation (EU) 2024/1689; Moffatt v Air Canada, 2024 BCCRT 149; ISO CG 40 47 / market coverage notes; SOC 2 AI practice write-ups 2025–26.",
        size=10,
        color=MUTED,
    )
    footer(s, 3)

    # 4 Why now
    s = blank_slide(prs)
    title_block(s, "Agents shipped. Evidence didn't.")
    card(s, Inches(0.8), Inches(1.35), Inches(5.7), Inches(4.5))
    add_text(s, Inches(1.1), Inches(1.6), Inches(5.1), Inches(0.35), "What got built", size=16, bold=True, color=ACCENT)
    add_paras(
        s,
        Inches(1.1),
        Inches(2.15),
        Inches(5.1),
        Inches(3.3),
        [
            "LangGraph, CrewAI, OpenAI Agents SDK, MCP, n8n — all optimized to make agents work.",
            "None of them give you a signed, exportable story of what ran.",
            "So people bolt on logs after the fact and hope it holds.",
        ],
        size=15,
        color=TEXT,
        spacing=14,
    )
    card(s, Inches(6.8), Inches(1.35), Inches(5.7), Inches(4.5))
    add_text(s, Inches(7.1), Inches(1.6), Inches(5.1), Inches(0.35), "What caught up", size=16, bold=True, color=ACCENT)
    add_paras(
        s,
        Inches(7.1),
        Inches(2.15),
        Inches(5.1),
        Inches(3.3),
        [
            "EU AI Act logging and oversight for high-risk systems.",
            "SOC 2 reviews starting to ask agent-shaped questions.",
            "NIST AI RMF — provenance language is no longer optional-sounding.",
        ],
        size=15,
        color=TEXT,
        spacing=14,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(6.15),
        Inches(11.5),
        Inches(0.45),
        "We didn't start Salanor because agents are exciting. We started it because the receipts have to exist before the ugly incident.",
        size=14,
        color=MUTED,
    )
    footer(s, 4)

    # 5 Solution — SDK + connectors
    s = blank_slide(prs)
    title_block(s, "Instrument once. Get a signed record.")
    add_text(
        s,
        Inches(0.8),
        Inches(1.15),
        Inches(11.5),
        Inches(0.55),
        "Same job whether you write code or run workflows: capture → sign → enforce → chain → export.\nEntry points differ. The ledger doesn't.",
        size=15,
        color=MUTED,
    )
    card(s, Inches(0.8), Inches(1.9), Inches(5.7), Inches(3.6))
    add_text(s, Inches(1.15), Inches(2.15), Inches(5.1), Inches(0.35), "In code — SDK", size=17, bold=True, color=ACCENT)
    add_paras(
        s,
        Inches(1.15),
        Inches(2.7),
        Inches(5.1),
        Inches(2.5),
        [
            "Wrap LangGraph, CrewAI, OpenAI Agents SDK, MCP tool calls.",
            "A few lines. Not a rewrite of your agent logic.",
        ],
        size=15,
        color=TEXT,
        spacing=12,
    )
    card(s, Inches(6.8), Inches(1.9), Inches(5.7), Inches(3.6))
    add_text(s, Inches(7.15), Inches(2.15), Inches(5.1), Inches(0.35), "In workflows — connectors", size=17, bold=True, color=ACCENT)
    add_paras(
        s,
        Inches(7.15),
        Inches(2.7),
        Inches(5.1),
        Inches(2.5),
        [
            "n8n node first (shipping / in active use on our side).",
            "Same signed events into Aegis — for teams that don't live in an IDE.",
            "More connectors as partners ask for them.",
        ],
        size=15,
        color=TEXT,
        spacing=12,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(5.85),
        Inches(11.5),
        Inches(0.6),
        "1 Capture   ·   2 Sign (your KMS)   ·   3 Enforce (<5ms p50)   ·   4 Chain   ·   5 Export",
        size=15,
        bold=True,
        color=WHITE,
    )
    footer(s, 5)

    # 6 Flow
    s = blank_slide(prs)
    title_block(s, "From action to something you can take to audit")
    flow = [
        ("1. Instrument", "SDK or n8n. Tool calls and LLM turns recorded as they happen."),
        ("2. Sign & enforce", "Signed at the point of action. Policy before it runs. Keys stay with you."),
        ("3. Ledger", "Hash-chained. Merkle roots so others can verify without your keys."),
        ("4. Export", "Rebuild a trace fast. Stream to Splunk, Datadog, or Sentinel."),
    ]
    for i, (t, b) in enumerate(flow):
        x = Inches(0.8) + Inches(3.05) * i
        card(s, x, Inches(1.5), Inches(2.9), Inches(4.3))
        add_text(s, x + Inches(0.2), Inches(1.8), Inches(2.5), Inches(0.7), t, size=16, bold=True, color=ACCENT)
        add_text(s, x + Inches(0.2), Inches(2.7), Inches(2.5), Inches(2.6), b, size=14, color=TEXT)
    footer(s, 6)

    # 7 Console
    s = blank_slide(prs)
    title_block(s, "What you see in the console")
    views = [
        ("Trace", "Every step in order, signature on each one."),
        ("Policy", "Allowed, blocked, and why — before the postmortem."),
        ("Export", "A bundle an auditor can use. Not a raw dump."),
        ("Verify", "Check a trace wasn't altered — without your keys."),
    ]
    for i, (t, b) in enumerate(views):
        row, col = divmod(i, 2)
        x = Inches(0.8) + Inches(6.1) * col
        y = Inches(1.4) + Inches(2.2) * row
        card(s, x, y, Inches(5.8), Inches(1.95))
        add_text(s, x + Inches(0.35), y + Inches(0.35), Inches(5.1), Inches(0.4), t, size=20, bold=True, color=WHITE)
        add_text(s, x + Inches(0.35), y + Inches(0.95), Inches(5.1), Inches(0.7), b, size=15, color=MUTED)
    add_text(s, Inches(0.8), Inches(6.15), Inches(11), Inches(0.35), "docs.salanor.com   ·   app.salanor.com", size=14, color=ACCENT)
    footer(s, 7)

    # 8 Open standard
    s = blank_slide(prs)
    title_block(s, "The format is open. We sell the infrastructure.")
    blocks = [
        ("APS-1", "Open provenance standard for agent events. CC BY 4.0."),
        ("Free verifier", "MIT CLI. Audit a trace without Salanor being online."),
        ("What we sell", "Managed control plane: hosting, policy, storage, exports."),
        ("BYOK", "Your keys. Ed25519. We can't rewrite your history. That's the point."),
    ]
    for i, (t, b) in enumerate(blocks):
        y = Inches(1.35) + Inches(1.15) * i
        card(s, Inches(0.8), y, Inches(11.7), Inches(1.0))
        add_text(s, Inches(1.15), y + Inches(0.28), Inches(2.8), Inches(0.4), t, size=17, bold=True, color=ACCENT)
        add_text(s, Inches(4.2), y + Inches(0.28), Inches(7.9), Inches(0.5), b, size=15, color=TEXT)
    footer(s, 8)

    # 9 Buyers + insurance note
    s = blank_slide(prs)
    title_block(s, "Who this is for")
    buyers = [
        ("Banks & fintech", "Model risk before agents touch underwriting or fraud."),
        ("Health systems", "Governance that won't approve autonomous triage without a reviewable trail."),
        ("Enterprise platforms", "You ship agents to customers. They'll ask what your product actually did."),
    ]
    for i, (t, b) in enumerate(buyers):
        x = Inches(0.8) + Inches(4.0) * i
        card(s, x, Inches(1.35), Inches(3.7), Inches(2.7))
        add_text(s, x + Inches(0.25), Inches(1.6), Inches(3.2), Inches(0.45), t, size=17, bold=True, color=WHITE)
        add_text(s, x + Inches(0.25), Inches(2.25), Inches(3.2), Inches(1.4), b, size=14, color=MUTED)
    card(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(1.7))
    add_text(s, Inches(1.15), Inches(4.5), Inches(11), Inches(0.35), "Insurance — later, on purpose", size=15, bold=True, color=ACCENT)
    add_text(
        s,
        Inches(1.15),
        Inches(5.0),
        Inches(11),
        Inches(0.8),
        "Carriers will need evidence before they price agent liability. That's Aether (risk / telemetry on top of signed Aegis patterns) — research, 2027.\nAegis first: the people running agents. Insurers second: once there's a trail worth underwriting.",
        size=14,
        color=TEXT,
    )
    footer(s, 9)

    # 10 Landscape
    s = blank_slide(prs)
    title_block(s, "Where we sit")
    add_text(
        s,
        Inches(0.8),
        Inches(1.15),
        Inches(11.5),
        Inches(0.4),
        "We're not another observability dashboard. Those help engineers. They aren't signed provenance for audit.",
        size=14,
        color=MUTED,
    )
    rows = [
        ("Agent frameworks / n8n", "Run agents and workflows. Don't produce evidence."),
        ("LLM observability", "Debug traces. Editable. Built for eng, not auditors."),
        ("SIEM / generic logs", "Can store events. Don't sign, chain, or export agent-shaped bundles."),
        ("Aegis", "Signed, hash-chained provenance + policy in path + exports. SDK and connectors."),
    ]
    for i, (t, b) in enumerate(rows):
        y = Inches(1.75) + Inches(1.05) * i
        card(s, Inches(0.8), y, Inches(11.7), Inches(0.9))
        add_text(s, Inches(1.15), y + Inches(0.25), Inches(3.5), Inches(0.4), t, size=15, bold=True, color=ACCENT if i == 3 else WHITE)
        add_text(s, Inches(4.8), y + Inches(0.25), Inches(7.3), Inches(0.5), b, size=14, color=TEXT)
    footer(s, 10)

    # 11 Where we are
    s = blank_slide(prs)
    title_block(s, "Where we are")
    meta = [
        ("Stage", "Prototype → design partner pilots"),
        ("HQ", "Kigali, Rwanda (Salanor Ltd)"),
        ("GTM", "Sales-led. One conversation at a time."),
        ("Aegis GA", "Target Q4 2026"),
    ]
    for i, (t, b) in enumerate(meta):
        x = Inches(0.8) + Inches(3.05) * i
        card(s, x, Inches(1.35), Inches(2.9), Inches(1.7))
        add_text(s, x + Inches(0.2), Inches(1.55), Inches(2.5), Inches(0.3), t, size=12, bold=True, color=ACCENT)
        add_text(s, x + Inches(0.2), Inches(2.0), Inches(2.5), Inches(0.8), b, size=14, color=TEXT)
    add_text(s, Inches(0.8), Inches(3.4), Inches(11), Inches(0.35), "In flight", size=16, bold=True, color=WHITE)
    add_paras(
        s,
        Inches(0.8),
        Inches(3.9),
        Inches(11.5),
        Inches(2.4),
        [
            "· Console — traces, policy, exports",
            "· APS-1 ingest + policy in the hot path",
            "· Compliance export bundles",
            "· n8n connector alongside the SDK",
            "· Spec + verifier already public",
        ],
        size=15,
        color=MUTED,
        spacing=7,
    )
    footer(s, 11)

    # 12 Business
    s = blank_slide(prs)
    title_block(s, "How we charge")
    card(s, Inches(0.8), Inches(1.35), Inches(5.7), Inches(4.0))
    add_text(s, Inches(1.15), Inches(1.65), Inches(5.1), Inches(0.35), "SaaS control plane", size=17, bold=True, color=ACCENT)
    add_paras(
        s,
        Inches(1.15),
        Inches(2.2),
        Inches(5.1),
        Inches(2.8),
        [
            "Per organization.",
            "Tiers by event volume and team size.",
            "Enterprise: SSO, region, retention, SLA.",
        ],
        size=15,
        color=TEXT,
        spacing=12,
    )
    card(s, Inches(6.8), Inches(1.35), Inches(5.7), Inches(4.0))
    add_text(s, Inches(7.15), Inches(1.65), Inches(5.1), Inches(0.35), "What's free", size=17, bold=True, color=ACCENT)
    add_paras(
        s,
        Inches(7.15),
        Inches(2.2),
        Inches(5.1),
        Inches(2.8),
        [
            "APS-1 stays free.",
            "Verifier stays free.",
            "We make money running Aegis — not locking the format.",
        ],
        size=15,
        color=TEXT,
        spacing=12,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(5.7),
        Inches(11.5),
        Inches(0.5),
        "No self-serve billing yet. Early customers are ops-led while we get pricing right. No fake ARR projections on this slide.",
        size=14,
        color=MUTED,
    )
    footer(s, 12)

    # 13 Team + ask
    s = blank_slide(prs)
    title_block(s, "Team & ask")
    card(s, Inches(0.8), Inches(1.3), Inches(5.7), Inches(3.5))
    add_text(s, Inches(1.15), Inches(1.55), Inches(5.1), Inches(0.35), "Landry Bougang Fotso", size=17, bold=True, color=WHITE)
    add_text(s, Inches(1.15), Inches(2.0), Inches(5.1), Inches(0.3), "Founder & CEO", size=13, color=ACCENT)
    add_text(
        s,
        Inches(1.15),
        Inches(2.5),
        Inches(5.1),
        Inches(2.0),
        "Building the layer that lets high-stakes agent decisions be explained later — not just executed.\n\nPreviously Aether (air quality) via Rwanda's Plug-In Ventures; pivoted here after discovery.",
        size=14,
        color=MUTED,
    )
    card(s, Inches(6.8), Inches(1.3), Inches(5.7), Inches(3.5))
    add_text(s, Inches(7.15), Inches(1.55), Inches(5.1), Inches(0.4), "A few design partners", size=16, bold=True, color=WHITE)
    add_paras(
        s,
        Inches(7.15),
        Inches(2.15),
        Inches(5.1),
        Inches(2.3),
        [
            "1. Agents where being wrong is expensive.",
            "2. You need signed provenance before scaling past pilot.",
            "3. Compliance/legal will help shape exports they'd accept.",
        ],
        size=14,
        color=TEXT,
        spacing=10,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(5.2),
        Inches(11.5),
        Inches(0.9),
        "partners@salanor.com\nsalanor.com  ·  docs.salanor.com  ·  salanor.com/spec",
        size=16,
        color=ACCENT,
    )
    add_text(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.35), "Aegis — provenance you can take to audit.", size=14, color=MUTED)
    footer(s, 13)

    prs.save(OUT)
    prs.save(OUT_REPO)
    print(f"Wrote {OUT}")
    print(f"Wrote {OUT_REPO}")


if __name__ == "__main__":
    build()
